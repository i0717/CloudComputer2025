# cli.py - 命令行界面
import os
import sys
import argparse
import json
from pathlib import Path
from typing import Optional
from dotenv import load_dotenv

load_dotenv()


# 动态导入，避免依赖问题
def import_components():
    """动态导入组件"""
    try:
        # 导入PPT解析器
        sys.path.append('.')
        from ppt_simple import SimplePPTParser
        from llm_direct import DirectLLMClient
        return SimplePPTParser, DirectLLMClient
    except ImportError as e:
        print(f"❌ 导入失败: {e}")
        print("请确保已安装所有依赖")
        sys.exit(1)


def main():
    """命令行主函数"""
    parser = argparse.ArgumentParser(
        description='PPT内容扩展智能体 - 命令行版本',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  %(prog)s --file presentation.pptx           # 处理单个文件
  %(prog)s --file slides.pptx --output results  # 指定输出目录
  %(prog)s --file notes.pptx --slides 1,3,5    # 只处理特定幻灯片
  %(prog)s --test                            # 运行测试
        """
    )

    parser.add_argument('--file', type=str, help='PPT文件路径')
    parser.add_argument('--output', type=str, default='output', help='输出目录')
    parser.add_argument('--slides', type=str, help='指定幻灯片编号（逗号分隔，如: 1,3,5）')
    parser.add_argument('--test', action='store_true', help='运行测试')
    parser.add_argument('--verbose', '-v', action='store_true', help='详细输出')

    args = parser.parse_args()

    if args.test:
        run_tests()
        return

    if not args.file:
        print("❌ 错误: 请提供PPT文件路径")
        parser.print_help()
        sys.exit(1)

    # 处理文件
    process_file(args.file, args.output, args.slides, args.verbose)


def run_tests():
    """运行测试"""
    print("🧪 运行系统测试...")

    # 检查API密钥
    api_key = os.getenv("SILICONFLOW_API_KEY")
    if not api_key:
        print("❌ 错误: 未设置 SILICONFLOW_API_KEY")
        print("请在 .env 文件中添加: SILICONFLOW_API_KEY=你的API密钥")
        return

    print(f"✅ API密钥已设置: {api_key[:10]}...")

    # 测试PPT解析
    try:
        from pptx import Presentation

        # 创建测试PPT
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "测试幻灯片"

        test_file = "test_temp.pptx"
        prs.save(test_file)

        print("✅ PPT创建测试通过")

        # 清理
        if os.path.exists(test_file):
            os.remove(test_file)

    except Exception as e:
        print(f"❌ PPT测试失败: {e}")
        return

    print("✅ 所有测试通过！")


def process_file(file_path: str, output_dir: str, slides_spec: Optional[str], verbose: bool = False):
    """处理PPT文件"""
    print(f"📄 处理文件: {file_path}")

    # 检查文件是否存在
    if not os.path.exists(file_path):
        print(f"❌ 文件不存在: {file_path}")
        sys.exit(1)

    # 检查文件类型
    if not file_path.lower().endswith(('.pptx', '.ppt')):
        print("❌ 仅支持PPT/PPTX文件")
        sys.exit(1)

    try:
        # 导入组件
        SimplePPTParser, DirectLLMClient = import_components()

        # 创建输出目录
        output_path = Path(output_dir)
        output_path.mkdir(exist_ok=True)

        # 1. 解析PPT
        print("🔍 解析PPT文件...")
        parser = SimplePPTParser()
        ppt_data = parser.parse(file_path)

        print(f"✅ 解析完成！找到 {ppt_data['total_slides']} 张幻灯片")

        if verbose:
            print("\n幻灯片预览:")
            for i, slide in enumerate(ppt_data['slides'][:3]):  # 只显示前3张
                print(f"  {i + 1}. {slide['title'][:30] or '无标题'}...")

        # 2. 确定要处理的幻灯片
        slides_to_process = []
        if slides_spec:
            # 解析幻灯片编号
            slide_numbers = []
            for num in slides_spec.split(','):
                try:
                    slide_numbers.append(int(num.strip()) - 1)  # 转换为0-based索引
                except ValueError:
                    print(f"❌ 无效的幻灯片编号: {num}")

            slides_to_process = [slide for i, slide in enumerate(ppt_data['slides'])
                                 if i in slide_numbers]
            print(f"📋 处理指定幻灯片: {slides_spec}")
        else:
            # 默认处理前5张或全部（如果少于5张）
            max_slides = min(5, len(ppt_data['slides']))
            slides_to_process = ppt_data['slides'][:max_slides]
            print(f"📋 处理前 {max_slides} 张幻灯片")

        if not slides_to_process:
            print("❌ 没有找到可处理的幻灯片")
            return

        # 3. 初始化LLM客户端
        print("🤖 初始化AI扩展器...")
        llm_client = DirectLLMClient()

        # 4. 扩展内容
        print("🧠 正在扩展内容（可能需要几分钟）...")

        expanded_results = []
        for i, slide in enumerate(slides_to_process):
            slide_num = slide['number'] + 1
            print(f"  正在处理幻灯片 {slide_num}/{len(slides_to_process)}: {slide['title'][:30] or '无标题'}...")

            # 准备扩展数据
            slide_for_expansion = {
                "title": slide['title'],
                "content": slide['text'][:3]  # 只取前3条内容
            }

            # 调用扩展
            try:
                result = llm_client.expand_slide(slide['title'], slide['text'])
                expanded_results.append({
                    "slide_number": slide_num,
                    "original": slide,
                    "expanded": result
                })

                if verbose:
                    print(f"    ✅ 扩展成功")
            except Exception as e:
                print(f"    ⚠️ 扩展失败: {e}")
                expanded_results.append({
                    "slide_number": slide_num,
                    "original": slide,
                    "error": str(e)
                })

        # 5. 保存结果
        print("💾 保存结果...")

        # 生成文件名
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        base_name = Path(file_path).stem

        # 保存JSON
        json_file = output_path / f"{base_name}_expanded_{timestamp}.json"
        with open(json_file, 'w', encoding='utf-8') as f:
            json.dump({
                "source_file": file_path,
                "processed_at": datetime.now().isoformat(),
                "total_slides": len(slides_to_process),
                "results": expanded_results
            }, f, ensure_ascii=False, indent=2)

        # 保存Markdown
        md_file = output_path / f"{base_name}_expanded_{timestamp}.md"
        save_as_markdown(file_path, expanded_results, md_file)

        print(f"\n✅ 处理完成！")
        print(f"📊 处理了 {len(expanded_results)} 张幻灯片")
        print(f"💾 结果已保存到:")
        print(f"   JSON: {json_file}")
        print(f"   Markdown: {md_file}")

        # 显示预览
        if expanded_results:
            print("\n📝 内容预览:")
            for result in expanded_results[:2]:  # 只显示前2个
                if 'expanded' in result:
                    content = result['expanded'].get('expanded_content', '')
                    if content:
                        print(f"幻灯片 {result['slide_number']}:")
                        print(f"  {content[:100]}...")
                        print()

    except Exception as e:
        print(f"❌ 处理失败: {e}")
        sys.exit(1)


def save_as_markdown(file_path: str, results: list, output_file: Path):
    """保存为Markdown格式"""
    from datetime import datetime

    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(f"# PPT内容扩展结果\n\n")
        f.write(f"**源文件**: {file_path}\n")
        f.write(f"**处理时间**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        f.write(f"**总幻灯片数**: {len(results)}\n\n")

        f.write("## 扩展内容\n\n")

        for result in results:
            slide_num = result['slide_number']
            original = result['original']

            f.write(f"### 幻灯片 {slide_num}: {original['title'] or '无标题'}\n\n")

            f.write("**原始内容**:\n")
            for item in original['text']:
                f.write(f"- {item}\n")
            f.write("\n")

            if 'expanded' in result:
                expanded = result['expanded']
                f.write("**扩展内容**:\n")
                f.write(expanded.get('expanded_content', '扩展失败') + "\n\n")

                sections = expanded.get('sections', {})
                if any(sections.values()):
                    f.write("**结构化内容**:\n")
                    for section, content in sections.items():
                        if content.strip():
                            f.write(f"#### {section}\n")
                            f.write(content + "\n\n")
            elif 'error' in result:
                f.write(f"**扩展失败**: {result['error']}\n\n")

            f.write("---\n\n")


if __name__ == "__main__":
    main()