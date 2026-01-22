import os
import logging
from typing import Dict, List, Optional, Any, Tuple
from pathlib import Path
import pptx
from pptx import Presentation
from pydantic import BaseModel
import json
from datetime import datetime
import re

logger = logging.getLogger(__name__)


class SlideContent(BaseModel):
    """幻灯片内容模型"""
    slide_number: int
    title: str
    content: List[str]
    bullet_points: List[str]
    images: List[str]
    notes: str = ""
    level: int = 1
    slide_type: str = "content"
    hierarchical_path: List[str] = []


class SlideStructure(BaseModel):
    """幻灯片结构分析结果"""
    slide_number: int
    title: str
    content_type: str
    hierarchical_level: int
    parent_titles: List[str] = []
    content_elements: List[Dict[str, Any]] = []
    has_images: bool = False
    has_tables: bool = False
    has_code: bool = False
    is_title_page: bool = False
    is_toc: bool = False
    is_empty: bool = False
    is_end_section: bool = False


class PPTMetadata(BaseModel):
    """PPT元数据"""
    filename: str
    total_slides: int
    author: Optional[str] = None
    created_date: Optional[str] = None
    modified_date:Optional[str] = None


class PPTStructure(BaseModel):
    """PPT结构"""
    metadata: PPTMetadata
    slides: List[SlideContent]
    outline: List[str]
    keywords: List[str]
    hierarchical_structure: List[SlideStructure] = []


class PPTParser:
    """PPT解析器"""

    def __init__(self):
        self.slides = []
        self.outline = []
        self.hierarchical_structure = []
        # 目录关键词 - 只要有这些词汇就认为是目录
        self.toc_keywords = [
            # 中文目录关键词
            "目录", "大纲", "议程", "日程", "内容提要", "内容概览",
            "主要议题", "课程大纲", "教学大纲", "报告大纲", "纲要",
            # 英文目录关键词
            "table of contents", "contents", "agenda", "schedule",
            "outline", "index", "syllabus", "curriculum", "program"
        ]
        
        # 章节标题关键词和模式
        self.chapter_patterns = [
            # 数字开头 + 顿号/点号 + 内容
            r'^\d+[、\.]\s*.+$',  # 1、xxx 或 1.xxx
            r'^[一二三四五六七八九十]+[、\.]\s*.+$',  # 一、xxx 或 一.xxx
            # 第X章格式
            r'^第[一二三四五六七八九十\d]+章[：:]?\s*.+$',  # 第1章：xxx
            r'^第[一二三四五六七八九十\d]+章节[：:]?\s*.+$',  # 第1章节：xxx
            # Chapter/Part格式
            r'^Chapter\s+\d+\s*[:：]?\s*.+$',  # Chapter 1: xxx
            r'^Part\s+\d+\s*[:：]?\s*.+$',  # Part 1: xxx
            # 单元/模块格式
            r'^单元[一二三四五六七八九十\d]+\s*[:：]?\s*.+$',  # 单元1：xxx
            r'^模块[一二三四五六七八九十\d]+\s*[:：]?\s*.+$',  # 模块1：xxx
            # 字母开头
            r'^[A-Z][、\.]\s*.+$',  # A、xxx 或 A.xxx
            r'^[A-Z]\s+.+$',  # A xxx
            # 带括号的数字
            r'^\(\d+\)\s*.+$',  # (1) xxx
            r'^（[一二三四五六七八九十\d]+）\s*.+$',  # （一）xxx
            # 其他常见格式
            r'^[一二三四五六七八九十\d]+[\.、]\s*.+$',  # 通用格式
        ]
        
        # 章节标题关键词 - 更全面的列表
        self.chapter_keywords = [
            "章", "章节", "节", "部分", "篇", "单元", "模块", "课", "讲", "讲次",
            "chapter", "part", "section", "unit", "module", "lesson", "lecture"
        ]
        
        # 小节标题模式
        self.section_patterns = [
            # 多级数字编号
            r'^\d+\.\d+\.\d+\s*.+$',  # 1.1.1 xxx
            r'^\d+\.\d+\s*.+$',  # 1.1 xxx
            r'^[一二三四五六七八九十]\.\d+\s*.+$',  # 一.1 xxx
            # 带括号的次级编号
            r'^\d+[\.、]\d+[\.、]\s*.+$',  # 1.1. xxx
            r'^[①②③④⑤⑥⑦⑧⑨⑩]\s*.+$',  # ① xxx
            # 小写字母开头
            r'^[a-z][、\.]\s*.+$',  # a、xxx 或 a.xxx
            r'^[a-z]\s+.+$',  # a xxx
            # 带括号的次级
            r'^\([a-z]\)\s*.+$',  # (a) xxx
            r'^（[一二三四五六七八九十]）\s*.+$',  # （一）xxx
            # 其他次级格式
            r'^\d+[\.、]\d+\s*.+$',  # 1.1 xxx
        ]
        
        # 小节标题关键词
        self.section_keywords = [
            "小节", "小标题", "标题", "要点", "主题", "话题", "子项", "子标题",
            "subsection", "subheading", "point", "topic", "subpoint", "subtitle"
        ]
        
        # 结尾页关键词
        self.end_keywords = [
            "结束", "完", "the end", "谢谢", "thank you", "q&a", "问答", 
            "讨论", "问题", "谢谢观看", "感谢聆听", "感谢倾听", "欢迎提问", 
            "any questions", "thanks", "感谢", "结语", "总结", "回顾",
            "future work", "conclusion", "ending", "finish", "done",
            "references", "参考文献", "致谢", "acknowledgement", "再见",
            "goodbye", "the end", "questions?", "q & a", "讨论与问答"
        ]
        
        # 摘要/总结关键词
        self.summary_keywords = [
            "总结", "小结", "要点回顾", "主要内容", "摘要", "概括",
            "summary", "conclusion", "key points", "takeaways", "recap",
            "概述", "概要", "核心内容", "重点", "要旨"
        ]
        
        # 目录页标记
        self.toc_found = False
        self.toc_slides = []  # 存储连续的目录页

    def parse_pptx(self, file_path: str) -> PPTStructure:
        """解析PPT文件"""
        try:
            logger.info(f"开始解析PPT文件: {file_path}")
            
            # 重置目录页标记
            self.toc_found = False
            self.toc_slides = []

            prs = Presentation(file_path)

            metadata = self._extract_metadata(prs, file_path)

            slides_content = []
            for i, slide in enumerate(prs.slides):
                slide_content = self._parse_slide_enhanced(slide, i)
                slides_content.append(slide_content)

            outline = self._extract_outline(slides_content)

            keywords = self._extract_keywords(slides_content)

            hierarchical_structure = self._analyze_hierarchical_structure(slides_content)

            structure = PPTStructure(
                metadata=metadata,
                slides=slides_content,
                outline=outline,
                keywords=keywords,
                hierarchical_structure=hierarchical_structure
            )

            logger.info(f"PPT解析完成: {len(slides_content)} 张幻灯片")
            return structure

        except Exception as e:
            logger.error(f"PPT解析失败: {e}")
            raise

    def _clean_text(self, text: str) -> str:
        """清理文本中的无效字符"""
        if not text:
            return ""
        
        # 移除控制字符（除了制表符、换行符、回车符）
        # 保留制表符(\t)、换行符(\n)、回车符(\r)
        cleaned = []
        for char in text:
            ord_val = ord(char)
            # 允许的字符：制表符(9)、换行符(10)、回车符(13)
            # 以及常规的可打印字符（32-126，中文字符等）
            if ord_val == 9 or ord_val == 10 or ord_val == 13:
                cleaned.append(char)
            elif ord_val >= 32 and ord_val <= 126:
                cleaned.append(char)
            elif ord_val >= 0x4e00 and ord_val <= 0x9fff:  # 中文字符范围
                cleaned.append(char)
            elif ord_val >= 0x3040 and ord_val <= 0x309F:  # 日文平假名
                cleaned.append(char)
            elif ord_val >= 0x30A0 and ord_val <= 0x30FF:  # 日文片假名
                cleaned.append(char)
            elif ord_val >= 0xAC00 and ord_val <= 0xD7AF:  # 韩文字符
                cleaned.append(char)
            elif ord_val >= 0xFF00 and ord_val <= 0xFFEF:  # 全角字符
                cleaned.append(char)
            elif char in ['•', '◦', '▪', '‣', '⁃', '∙', '○', '◉', '◎', '✓', '✔', '→', '➔', '➜', '➤']:
                cleaned.append(char)  # 保留项目符号
            elif char in ['、', '。', '，', '；', '：', '？', '！', '「', '」', '『', '』', '（', '）', '【', '】', '《', '》']:
                cleaned.append(char)  # 保留中文标点
            elif char in ['·', '•', '…', '—', '～', '＇', '＂', '＃', '＄', '％', '＆', '＇', '（', '）', '＊', '＋', '，', '－', '．', '／']:
                cleaned.append(char)  # 保留其他特殊字符
            else:
                # 替换其他控制字符为空格
                cleaned.append(' ')
        
        cleaned_text = ''.join(cleaned)
        # 合并多个空格
        cleaned_text = re.sub(r'\s+', ' ', cleaned_text)
        return cleaned_text.strip()

    def _extract_metadata(self, presentation, file_path: str) -> PPTMetadata:
        """提取元数据"""
        try:
            core_props = presentation.core_properties

            return PPTMetadata(
                filename=Path(file_path).name,
                total_slides=len(presentation.slides),
                author=core_props.author,
                created_date=core_props.created.isoformat() if core_props.created else None,
                modified_date=core_props.modified.isoformat() if core_props.modified else None
            )
        except:
            return PPTMetadata(
                filename=Path(file_path).name,
                total_slides=len(presentation.slides)
            )

    def _parse_slide_enhanced(self, slide, slide_num: int) -> SlideContent:
        """增强版幻灯片解析"""
        title = ""
        content = []
        bullet_points = []
        
        # 收集所有文本，不在乎有几个对话框
        all_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                try:
                    text = shape.text
                    if text:
                        # 清理无效字符
                        text = self._clean_text(text)
                        if text:  # 清理后可能为空
                            all_texts.append(text)
                except Exception as e:
                    logger.warning(f"处理幻灯片 {slide_num} 的形状文本时出错: {e}")
                    continue
        
        # 检查每一段文本，寻找可能的标题
        potential_titles = []
        all_other_texts = []
        
        for text in all_texts:
            # 只有小于50个字的文本才考虑作为标题
            if len(text) <= 50:
                potential_titles.append(text)
            else:
                all_other_texts.append(text)
        
        # 如果有候选标题，选择第一个作为标题
        if potential_titles:
            title = potential_titles[0]
            # 剩下的候选标题也作为内容
            if len(potential_titles) > 1:
                all_other_texts.extend(potential_titles[1:])
        
        # 处理所有其他文本
        for text in all_other_texts:
            if self._is_bullet_point(text):
                bullet_points.append(text)
            else:
                content.append(text)

        images = []
        image_count = 0
        for i, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # 图片类型
                images.append(f"slide_{slide_num}_image_{image_count}")
                image_count += 1
            elif hasattr(shape, "image"):
                images.append(f"slide_{slide_num}_image_{image_count}")
                image_count += 1
            elif hasattr(shape, "fill") and hasattr(shape.fill, "type"):
                if shape.fill.type == 6:  # 图片填充
                    images.append(f"slide_{slide_num}_image_{image_count}")
                    image_count += 1

        if len(images) == 0 and slide.background and hasattr(slide.background, "fill"):
            if slide.background.fill.type == 6:
                images.append(f"slide_{slide_num}_background_image")

        # 清理备注文本
        notes = ""
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text:
                    notes = self._clean_text(notes_text)
            except Exception as e:
                logger.warning(f"处理幻灯片 {slide_num} 的备注时出错: {e}")

        level = self._determine_level(title, content, bullet_points)

        return SlideContent(
            slide_number=slide_num,
            title=title,
            content=content,
            bullet_points=bullet_points,
            images=images,
            notes=notes,
            level=level
        )

    def _is_bullet_point(self, text: str) -> bool:
        """检测是否为项目符号文本"""
        # 常见的项目符号字符
        bullet_chars = ["•", "◦", "▪", "‣", "⁃", "∙", "○", "◉", "◎", "✓", "✔", "→", "➔", "➜", "➤"]
        
        # 数字或字母开头的列表项
        number_patterns = [
            r'^\d+[\.、\)）]',      # 数字后接点、顿号、括号
            r'^[a-zA-Z][\.、\)）]',  # 字母后接点、顿号、括号
            r'^[一二三四五六七八九十][\.、\)）]',  # 中文数字
        ]
        
        # 检查是否以项目符号开头
        if any(text.startswith(bullet) or bullet in text[:3] for bullet in bullet_chars):
            return True
        
        # 检查是否匹配数字/字母列表模式
        for pattern in number_patterns:
            if re.match(pattern, text.strip()):
                return True
        
        # 检查是否有缩进（可能表示项目符号）
        if text.startswith('  ') or text.startswith('\t'):
            return True
            
        return False

    def _determine_level(self, title: str, content: List[str], bullet_points: List[str]) -> int:
        """确定内容层级 - 改进版"""
        all_text = title + " ".join(content) + " ".join(bullet_points)
        
        # 根据文本特征判断层级
        if self._is_chapter_title_simple(title):
            return 1
        elif self._is_section_title_simple(title):
            return 2
        elif self._looks_like_subsection(title, content):
            return 3
        else:
            return 4

    def _looks_like_main_title(self, title: str, content: List[str]) -> bool:
        """判断是否像主标题"""
        if not title:
            return False
            
        # 主标题通常较短，且可能包含冒号分隔
        title_clean = title.replace("：", ":").replace("——", ":")
        
        # 检查是否为标题页格式
        title_page_patterns = [
            r'^[^:]*:[^:]*$',  # 包含一个冒号
            r'^.*报告$',        # 以"报告"结尾
            r'^.*演讲$',        # 以"演讲"结尾
            r'^.*介绍$',        # 以"介绍"结尾
        ]
        
        for pattern in title_page_patterns:
            if re.match(pattern, title_clean):
                return True
        
        # 如果标题较短且内容很少或为空
        if len(title) < 30 and len(content) == 0:
            return True
            
        return False

    def _looks_like_section(self, title: str, content: List[str]) -> bool:
        """判断是否像章节"""
        if not title:
            return False
            
        # 章节通常有编号
        section_patterns = [
            r'^第[一二三四五六七八九十\d]+[章节部分]',
            r'^Chapter\s+\d+',
            r'^Part\s+\d+',
            r'^\d+\.\s+',
            r'^[一二三四五六七八九十]、',
        ]
        
        for pattern in section_patterns:
            if re.match(pattern, title.strip()):
                return True
                
        return False

    def _looks_like_subsection(self, title: str, content: List[str]) -> bool:
        """判断是否像小节"""
        if not title:
            return False
            
        # 小节通常有次级编号
        subsection_patterns = [
            r'^\d+\.\d+\.',
            r'^[一二三四五六七八九十]\.\d+',
            r'^\d+、',
            r'^\([一二三四五六七八九十\d]+\)',
            r'^[a-z]\.',
        ]
        
        for pattern in subsection_patterns:
            if re.match(pattern, title.strip()):
                return True
                
        return False

    def _extract_outline(self, slides: List[SlideContent]) -> List[str]:
        """提取大纲"""
        outline = []
        for slide in slides:
            if slide.level <= 3 and slide.title:
                indent = "  " * (slide.level - 1)
                outline.append(f"{indent}{slide.title}")
        return outline

    def _extract_keywords(self, slides: List[SlideContent]) -> List[str]:
        """提取关键词"""
        keywords = set()
        
        for slide in slides:
            # 从标题提取关键词
            if slide.title:
                # 清理标题
                clean_title = self._clean_text(slide.title)
                if clean_title:
                    # 分割标题，提取有意义的部分
                    parts = re.split(r'[：:;；,，\-—]', clean_title)
                    for part in parts:
                        part_clean = part.strip()
                        if 2 <= len(part_clean) <= 20 and not part_clean.isdigit():
                            keywords.add(part_clean)
            
            # 从内容中提取可能的关键词（较长的名词短语）
            for text in slide.content:
                clean_text = self._clean_text(text)
                if len(clean_text) > 10:
                    # 提取中文短语
                    chinese_phrases = re.findall(r'[\u4e00-\u9fa5]{2,6}', clean_text)
                    keywords.update(chinese_phrases)
        
        return list(keywords)[:15]  # 返回前15个关键词

    def _analyze_hierarchical_structure(self, slides: List[SlideContent]) -> List[SlideStructure]:
        """分析PPT的层级结构"""
        if not slides:
            return []
        
        structure_results = []
        
        # 处理第一页（通常是标题页）
        first_slide = slides[0]
        first_structure = self._classify_first_slide(first_slide)
        structure_results.append(first_structure)
        
        # 当前层级路径
        current_hierarchy = [first_structure.title] if first_structure.title != "主标题" else ["主标题"]
        
        # 处理后续幻灯片
        for i in range(1, len(slides)):
            slide = slides[i]
            
            # 检查是否为目录页 - 改进的检测，支持连续目录页
            is_toc = self._is_toc_slide_improved(slide, i, slides)
            
            if is_toc:
                structure = self._create_toc_structure(slide, current_hierarchy)
                structure_results.append(structure)
                # 更新层级路径
                if "目录" not in current_hierarchy and "大纲" not in current_hierarchy:
                    current_hierarchy.append("目录" if slide.title != "大纲" else "大纲")
                continue
            
            # 检查是否为结尾页 - 只检查最后一页
            is_end, end_type = self._detect_end_page(slide, i, len(slides))
            if is_end:
                structure = self._create_end_structure(slide, end_type)
                structure_results.append(structure)
                current_hierarchy = ["结尾部分"]
                continue
            
            # 分类幻灯片内容 - 按照要求的顺序，且一旦识别就确定，不再被覆盖
            content_type = self._classify_slide_content_strict(slide, i, len(slides))
            
            # 确定层级
            hierarchical_level = self._determine_hierarchical_level(
                slide, content_type, current_hierarchy, i
            )
            
            # 分析内容元素
            content_elements = self._analyze_content_elements(slide, content_type)
            
            # 创建结构对象
            structure = SlideStructure(
                slide_number=slide.slide_number,
                title=slide.title or f"幻灯片 {slide.slide_number + 1}",
                content_type=content_type,
                hierarchical_level=hierarchical_level,
                parent_titles=current_hierarchy.copy(),
                content_elements=content_elements,
                has_images=len(slide.images) > 0,
                has_tables=self._contains_tables(slide),
                has_code=self._contains_code(slide),
                is_title_page=False,
                is_toc=is_toc,
                is_empty=self._is_empty_slide(slide),
                is_end_section=is_end
            )
            structure_results.append(structure)
            
            # 更新层级路径（对于非目录、非结尾页）
            if not is_toc and not is_end:
                self._update_hierarchy_path(structure, current_hierarchy)
        
        return structure_results

    def _classify_slide_content_strict(self, slide: SlideContent, slide_num: int, total_slides: int) -> str:
        """严格的幻灯片内容类型分类 - 按照要求的顺序，且一旦识别就不再被覆盖"""
        
        # 0. 检查是否为空（这个最先检查）
        if self._is_empty_slide(slide):
            return "空白页"
        
        # 1. 检查是否为主标题 - 只有第一页可以识别为主标题
        if slide_num == 0:
            if self._is_title_page(slide):
                return "主标题"
        
        # 2. 检查是否为目录
        is_toc = self._is_toc_slide_improved_simple(slide, slide_num, total_slides)
        if is_toc:
            return "目录"
        
        # 3. 检查是否为结尾等
        is_end, end_type = self._detect_end_page_simple(slide, slide_num, total_slides)
        if is_end:
            if end_type in ["致谢", "参考文献", "问答"]:
                return end_type
            else:
                return "结尾页"
        
        # 4. 检查是否为章节标题 - 修改：只有当一页上的字总数少于50个字时才判断
        if slide.title and self._count_slide_text_chars(slide) <= 50:
            if self._is_chapter_title_simple(slide.title):
                return "章节标题"
        
        # 5. 检查是否为图片页 - 字数限制在10以内
        if self._is_image_page(slide):
            return "图片页"
        
        # 6. 检查是否为正文
        has_title = bool(slide.title and slide.title.strip())
        has_content = bool(slide.content or slide.bullet_points)
        
        if has_title or has_content:
            # 检查是否是小节标题 - 同样限制：只有当一页上的字总数少于50个字时才判断
            if slide.title and self._count_slide_text_chars(slide) <= 50:
                if self._is_section_title_simple(slide.title):
                    return "小节标题"
            
            return "正文"
        
        # 7. 默认返回空白页
        return "空白页"

    def _count_slide_text_chars(self, slide: SlideContent) -> int:
        """计算幻灯片中的总字符数（包括标题、内容和项目符号）"""
        total_text = ""
        
        if slide.title:
            total_text += slide.title
        
        for text in slide.content:
            total_text += text
        
        for bullet in slide.bullet_points:
            total_text += bullet
        
        # 只计算有意义的字符（排除空格和标点）
        chinese_chars = re.findall(r'[\u4e00-\u9fa5]', total_text)
        english_letters = re.findall(r'[a-zA-Z]', total_text)
        numbers = re.findall(r'\d', total_text)
        
        total_count = len(chinese_chars) + len(english_letters) + len(numbers)
        
        return total_count

    def _is_toc_slide_improved_simple(self, slide: SlideContent, slide_num: int, total_slides: int) -> bool:
        """简化版目录页检测"""
        title = slide.title or ""
        
        # 检查标题中是否包含目录关键词
        title_lower = title.lower()
        for keyword in self.toc_keywords:
            if keyword.lower() in title_lower:
                return True
        
        # 如果标题中没有，检查内容中是否有目录关键词
        all_text = " ".join(slide.content + slide.bullet_points)
        text_lower = all_text.lower()
        for keyword in self.toc_keywords:
            if keyword.lower() in text_lower:
                return True
        
        return False

    def _detect_end_page_simple(self, slide: SlideContent, slide_num: int, total_slides: int) -> Tuple[bool, str]:
        """简化版结尾页检测"""
        # 只检查最后一页
        if slide_num != total_slides - 1:
            return False, ""
        
        title = slide.title or ""
        all_text = " ".join(slide.content + slide.bullet_points)
        text_lower = all_text.lower()
        
        # 检查是否包含结尾关键词
        for keyword in self.end_keywords:
            if keyword in text_lower or keyword in title.lower():
                # 进一步判断结尾页类型
                if any(kw in text_lower for kw in ["谢谢", "感谢", "thank"]):
                    return True, "致谢"
                elif any(kw in text_lower for kw in ["reference", "参考文献"]):
                    return True, "参考文献"
                elif any(kw in text_lower for kw in ["q&a", "问答"]):
                    return True, "问答"
                else:
                    return True, "结尾页"
        
        return False, ""

    def _is_chapter_title_simple(self, title: str) -> bool:
        """简化版章节标题检测 - 只检查标题文本，支持多种数字格式"""
        if not title or len(title.strip()) == 0:  # 纯数字"5"可以，但空字符串不行
            return False
        
        title_clean = title.strip()
        
        # 排除通用标题
        if re.match(r'^幻灯片\s*\d+$', title_clean):
            return False
        
        if re.match(r'^Slide\s*\d+$', title_clean, re.IGNORECASE):
            return False
        
        # 1. 检查是否以数字或中文数字开头，后跟顿号或点号 - 支持多种格式
        simple_patterns = [
            r'^\d+[、\.]\s*.+$',  # 1、xxx 或 1.xxx（包括01、02等）
            r'^[一二三四五六七八九十]+[、\.]\s*.+$',  # 一、xxx 或 一.xxx
            r'^第[一二三四五六七八九十\d]+[章节].*$',  # 第1章 xxx
            r'^[A-Z][、\.]\s*.+$',  # A、xxx 或 A.xxx
            r'^\(\d+\)\s*.+$',  # (1) xxx
            r'^第[一二三四五六七八九十\d]+部分.*$',  # 第1部分 xxx
            r'^第[一二三四五六七八九十\d]+单元.*$',  # 第1单元 xxx
            r'^第[一二三四五六七八九十\d]+讲.*$',  # 第1讲 xxx
            r'^第[一二三四五六七八九十\d]+课.*$',  # 第1课 xxx
            r'^Chapter\s+\d+.*$',  # Chapter 1 xxx
            r'^Part\s+\d+.*$',  # Part 1 xxx
            r'^\d+\s+.*$',  # 1 xxx 或 01 xxx（数字后跟空格，然后是文字）
            r'^\d+$',  # 纯数字（如01、02、5、10等）- 这个模式应该能匹配"5"
        ]
        
        for pattern in simple_patterns:
            if re.match(pattern, title_clean):
                return True
        
        # 2. 检查是否以01、02等格式开头（数字后可能有空格）
        if re.match(r'^0\d+\s*.+$', title_clean):  # 01 xxx, 02 xxx等
            return True
        
        # 3. 检查是否以单个数字开头，后跟空格和文字（如"1 引言"、"2 方法"、"3 结果"）
        # 这个模式专门匹配数字（1-9）后跟空格，然后是至少一个非数字字符
        if re.match(r'^[1-9]\s+[^\d].+$', title_clean):
            return True
        
        # 4. 检查是否以单个数字开头，后跟空格和任何内容（包括数字）
        if re.match(r'^[1-9]\s+.+$', title_clean):  # 更宽松的模式
            return True
        
        # 5. 检查标题是否较短（章节标题通常较短）且包含章节关键词
        if len(title_clean) < 30:
            title_lower = title_clean.lower()
            for keyword in self.chapter_keywords:
                if keyword.lower() in title_lower:
                    return True
        
        # 6. 纯数字（如"5"）应该被认为是章节标题
        if title_clean.isdigit():
            return True
        
        return False

    def _is_section_title_simple(self, title: str) -> bool:
        """简化版小节标题检测"""
        if not title or len(title.strip()) < 1:  # 改为至少1个字符
            return False
        
        title_clean = title.strip()
        
        # 排除通用标题
        if re.match(r'^幻灯片\s*\d+$', title_clean):
            return False
        
        # 检查是否匹配小节标题模式
        simple_patterns = [
            r'^\d+\.\d+[、\.]?.+$',  # 1.1 xxx 或 1.1.xxx
            r'^[a-z][、\.]\s*.+$',  # a、xxx 或 a.xxx
            r'^\([a-z]\)\s*.+$',  # (a) xxx
            r'^[①②③④⑤⑥⑦⑧⑨⑩]\s*.+$',  # ① xxx
            r'^\d+\.\d+\.\d+[、\.]?.+$',  # 1.1.1 xxx
            r'^[一二三四五六七八九十]\.\d+[、\.]?.+$',  # 一.1 xxx
            r'^\(\d+\)\s*.+$',  # (1) xxx 也考虑为小节标题
            r'^\d+\.\d+[、\.]?.+$',  # 01.1 xxx 或 02.1 xxx
            r'^\d+\s+.*$',  # 01 xxx 或 02 xxx 或 1 xxx
            r'^\d+$',  # 纯数字如5
        ]
        
        for pattern in simple_patterns:
            if re.match(pattern, title_clean):
                return True
        
        # 检查是否包含小节关键词且标题较短
        if len(title_clean) < 25:
            title_lower = title_clean.lower()
            for keyword in self.section_keywords:
                if keyword.lower() in title_lower:
                    return True
        
        # 纯数字应该被认为是小节标题
        if title_clean.isdigit():
            return True
        
        return False

    def _is_toc_slide_improved(self, slide: SlideContent, slide_num: int, all_slides: List[SlideContent]) -> bool:
        """改进版目录页检测 - 支持连续目录页"""
        title = slide.title or ""
        
        # 检查标题中是否包含目录关键词
        title_lower = title.lower()
        has_toc_keyword = False
        for keyword in self.toc_keywords:
            if keyword.lower() in title_lower:
                has_toc_keyword = True
                break
        
        # 如果标题中没有，检查内容中是否有目录关键词
        if not has_toc_keyword:
            all_text = " ".join(slide.content + slide.bullet_points)
            text_lower = all_text.lower()
            for keyword in self.toc_keywords:
                if keyword.lower() in text_lower:
                    has_toc_keyword = True
                    break
        
        if not has_toc_keyword:
            return False
        
        # 如果是第一页目录页
        if not self.toc_found:
            self.toc_found = True
            self.toc_slides = [slide_num]
            return True
        
        # 如果是已经开始的连续目录页
        if self.toc_found:
            # 检查是否与上一个目录页连续
            if slide_num == self.toc_slides[-1] + 1:
                self.toc_slides.append(slide_num)
                return True
        
        return False

    def _classify_first_slide(self, slide: SlideContent) -> SlideStructure:
        """分类第一张幻灯片"""
        # 判断是否为标题页
        is_title_page = self._is_title_page(slide)
        
        if is_title_page:
            content_type = "主标题"
            title = slide.title or "主标题"
            if not title or title == f"幻灯片 1":
                if slide.content:
                    title = slide.content[0][:50]
                elif slide.bullet_points:
                    title = slide.bullet_points[0][:50]
        else:
            # 第一页如果不是标题页，就按照正常流程分类
            content_type = self._classify_slide_content_strict(slide, 0, 1)
            title = slide.title or f"幻灯片 1"
        
        return SlideStructure(
            slide_number=0,
            title=title,
            content_type=content_type,
            hierarchical_level=1,
            parent_titles=[],
            content_elements=self._analyze_content_elements(slide, content_type),
            has_images=len(slide.images) > 0,
            has_tables=self._contains_tables(slide),
            has_code=self._contains_code(slide),
            is_title_page=is_title_page,
            is_toc=False,
            is_empty=self._is_empty_slide(slide),
            is_end_section=False
        )

    def _is_title_page(self, slide: SlideContent) -> bool:
        """判断是否为标题页"""
        # 标题页通常有较短的标题，可能包含副标题
        title = slide.title or ""
        
        # 检查标题特征
        if len(title) > 50:  # 标题过长可能不是标题页
            return False
        
        # 检查是否包含标题页常见元素
        title_page_indicators = [
            "报告", "演讲", "汇报", "介绍", "presentation", "report",
            "培训", "讲座", "课程", "workshop", "seminar", "分享"
        ]
        
        title_lower = title.lower()
        for indicator in title_page_indicators:
            if indicator in title_lower:
                return True
        
        # 检查是否有副标题模式（冒号分隔）
        if ":" in title or "：" in title or "——" in title:
            return True
        
        # 如果幻灯片内容很少，可能是标题页
        total_text = len(title) + sum(len(c) for c in slide.content) + sum(len(b) for b in slide.bullet_points)
        if total_text < 100 and slide.slide_number == 0:
            return True
            
        return False

    def _create_toc_structure(self, slide: SlideContent, current_hierarchy: List[str]) -> SlideStructure:
        """创建目录页结构"""
        # 确定目录标题
        title = slide.title or "目录"
        if title == f"幻灯片 {slide.slide_number + 1}":
            # 尝试从内容中提取标题
            if slide.content:
                first_line = slide.content[0]
                if len(first_line) < 30:
                    title = first_line
            elif slide.bullet_points:
                first_bullet = slide.bullet_points[0]
                if len(first_bullet) < 30:
                    title = first_bullet
        
        return SlideStructure(
            slide_number=slide.slide_number,
            title=title,
            content_type="目录",
            hierarchical_level=len(current_hierarchy) + 1,
            parent_titles=current_hierarchy.copy(),
            content_elements=self._analyze_content_elements(slide, "目录"),
            has_images=len(slide.images) > 0,
            has_tables=self._contains_tables(slide),
            has_code=self._contains_code(slide),
            is_title_page=False,
            is_toc=True,
            is_empty=self._is_empty_slide(slide),
            is_end_section=False
        )

    def _detect_end_page(self, slide: SlideContent, slide_num: int, total_slides: int) -> Tuple[bool, str]:
        """检测结尾页 - 只检测最后一页"""
        # 只检查最后一页
        if slide_num != total_slides - 1:
            return False, ""
        
        title = slide.title or ""
        all_text = " ".join(slide.content + slide.bullet_points)
        text_lower = all_text.lower()
        
        # 检查是否包含结尾关键词
        for keyword in self.end_keywords:
            if keyword in text_lower or keyword in title.lower():
                # 进一步判断结尾页类型
                if any(kw in text_lower for kw in ["谢谢", "感谢", "thank"]):
                    return True, "致谢"
                elif any(kw in text_lower for kw in ["reference", "参考文献"]):
                    return True, "参考文献"
                elif any(kw in text_lower for kw in ["q&a", "问答"]):
                    return True, "问答"
                else:
                    return True, "结尾页"
        
        # 检查是否为摘要/总结（最后一页的摘要也算结尾）
        for keyword in self.summary_keywords:
            if keyword in text_lower or keyword in title.lower():
                return True, "结尾页"
        
        # 如果内容很少且有图片，可能是结尾页
        if len(slide.images) > 0 and self._count_text_chars(slide) < 10:
            return True, "结尾页"
        
        # 如果文本非常少
        if self._count_text_chars(slide) < 10:
            return True, "结尾页"
        
        return False, ""

    def _create_end_structure(self, slide: SlideContent, end_type: str) -> SlideStructure:
        """创建结尾页结构"""
        title = slide.title or end_type
        
        return SlideStructure(
            slide_number=slide.slide_number,
            title=title,
            content_type=end_type,
            hierarchical_level=1,
            parent_titles=["结尾部分"],
            content_elements=self._analyze_content_elements(slide, end_type),
            has_images=len(slide.images) > 0,
            has_tables=self._contains_tables(slide),
            has_code=self._contains_code(slide),
            is_title_page=False,
            is_toc=False,
            is_empty=self._is_empty_slide(slide),
            is_end_section=True
        )

    def _is_summary_page(self, title: str, all_text: str) -> bool:
        """判断是否为摘要/总结页"""
        text_lower = all_text.lower()
        title_lower = title.lower()
        
        for keyword in self.summary_keywords:
            if keyword in text_lower or keyword in title_lower:
                return True
        
        return False

    def _is_image_page(self, slide: SlideContent) -> bool:
        """检查是否为图片页 - 字数限制在10以内"""
        has_images = len(slide.images) > 0
        
        if not has_images:
            return False
        
        char_count = self._count_text_chars(slide)
        
        # 图片页：有图片且文字很少（10个字符以内）
        return char_count < 10

    def _contains_code(self, slide: SlideContent) -> bool:
        """检查是否包含代码"""
        all_text = " ".join(slide.content + slide.bullet_points)
        
        # 代码常见特征
        code_patterns = [
            r'def\s+\w+\(',  # 函数定义
            r'class\s+\w+',  # 类定义
            r'import\s+\w+',  # 导入语句
            r'print\(',  # 打印语句
            r'for\s+\w+\s+in',  # for循环
            r'if\s+\w+',  # if语句
            r'\{.*\}',  # 花括号
            r'=.*;',  # 分号结尾的赋值
        ]
        
        for pattern in code_patterns:
            if re.search(pattern, all_text):
                return True
        
        # 检查是否有缩进的代码块
        lines = all_text.split('\n')
        indented_lines = sum(1 for line in lines if line.startswith('    ') or line.startswith('\t'))
        if indented_lines >= 3:
            return True
        
        return False

    def _determine_hierarchical_level(self, slide: SlideContent, content_type: str, 
                                     current_hierarchy: List[str], slide_num: int) -> int:
        """确定层级深度"""
        if content_type in ["主标题", "目录"]:
            return 1
        
        if content_type in ["结尾页", "致谢", "参考文献", "问答"]:
            return 1
        
        if content_type == "章节标题":
            return 2
        
        if content_type == "小节标题":
            return 3
        
        if content_type == "图片页":
            return len(current_hierarchy) + 1
        
        # 其他类型（正文）根据当前层级确定
        return len(current_hierarchy) + 1

    def _count_text_chars(self, slide: SlideContent) -> int:
        """计算幻灯片中的总字符数"""
        total_text = ""
        
        if slide.title:
            total_text += slide.title
        
        for text in slide.content:
            total_text += text
        
        for bullet in slide.bullet_points:
            total_text += bullet
        
        # 只计算有意义的字符（排除空格和标点）
        chinese_chars = re.findall(r'[\u4e00-\u9fa5]', total_text)
        english_letters = re.findall(r'[a-zA-Z]', total_text)
        numbers = re.findall(r'\d', total_text)
        
        total_count = len(chinese_chars) + len(english_letters) + len(numbers)
        
        return total_count

    def _analyze_content_elements(self, slide: SlideContent, content_type: str) -> List[Dict[str, Any]]:
        """分析内容元素 - 简化版"""
        elements = []
        
        if slide.title:
            # 根据内容类型确定标题重要性
            if content_type in ["主标题", "章节标题", "图片页"]:
                importance = "high"
            elif content_type in ["目录", "结尾页", "致谢"]:
                importance = "medium"
            else:
                importance = "low"
                
            elements.append({
                "type": "标题",
                "content": slide.title,
                "format": "标题",
                "importance": importance
            })
        
        # 处理内容文本 - 统一归类为"正文"
        for i, text in enumerate(slide.content):
            # 简化文本元素分类
            if len(text) < 50:
                element_type = "短句"
            else:
                element_type = "正文段落"
            
            # 根据位置确定重要性
            if i == 0 and len(text) > 10:
                importance = "medium"
            else:
                importance = "low"
                
            elements.append({
                "type": element_type,
                "content": text,
                "format": "段落",
                "importance": importance
            })
        
        # 处理项目符号
        for bullet in slide.bullet_points:
            elements.append({
                "type": "项目符号",
                "content": bullet,
                "format": "列表项",
                "importance": "medium"
            })
        
        # 处理图片
        if slide.images:
            for i, image_ref in enumerate(slide.images):
                elements.append({
                    "type": "图片",
                    "content": f"图片 {i+1}",
                    "format": "图片",
                    "importance": "high" if content_type == "图片页" else "medium",
                    "image_ref": image_ref
                })
        
        # 处理备注
        if slide.notes:
            elements.append({
                "type": "备注",
                "content": slide.notes,
                "format": "备注",
                "importance": "low"
            })
        
        return elements

    def _contains_tables(self, slide: SlideContent) -> bool:
        """检查是否包含表格 - 移除表格检测"""
        # 根据要求，移除表格检测功能
        return False

    def _is_empty_slide(self, slide: SlideContent) -> bool:
        """检查是否为空幻灯片"""
        has_title = bool(slide.title and slide.title.strip())
        has_content = bool(slide.content or slide.bullet_points)
        has_images = bool(slide.images)
        
        return not (has_title or has_content or has_images)

    def _update_hierarchy_path(self, structure: SlideStructure, current_hierarchy: List[str]) -> None:
        """更新层级路径"""
        if structure.content_type in ["章节标题", "图片页"]:
            # 根据层级深度调整路径
            if structure.hierarchical_level <= len(current_hierarchy):
                current_hierarchy = current_hierarchy[:structure.hierarchical_level-1]
            
            if structure.title and structure.title not in current_hierarchy:
                current_hierarchy.append(structure.title)
        elif structure.content_type == "正文":
            # 正文通常不改变主层级，除非是新的章节
            pass

    def save_to_json(self, structure: PPTStructure, output_path: str):
        """保存解析结果到JSON文件"""
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(structure.dict(), f, ensure_ascii=False, indent=2)
        logger.info(f"解析结果已保存: {output_path}")

    def load_from_json(self, json_path: str) -> PPTStructure:
        """从JSON文件加载解析结果"""
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return PPTStructure(**data)


ppt_parser = PPTParser()
